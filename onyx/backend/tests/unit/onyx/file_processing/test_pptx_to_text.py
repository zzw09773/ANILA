import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from onyx.file_processing.extract_file_text import pptx_to_text


def _make_pptx_with_chart() -> io.BytesIO:
    """Create an in-memory pptx with one text slide and one chart slide."""
    prs = Presentation()

    # Slide 1: text only
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Introduction"
    slide1.placeholders[1].text = "This is the first slide."

    # Slide 2: chart
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3"]
    chart_data.add_series("Revenue", (100, 200, 300))
    slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(4),
        chart_data,
    )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _make_pptx_without_chart() -> io.BytesIO:
    """Create an in-memory pptx with a single text-only slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Hello World"
    slide.placeholders[1].text = "Some content here."

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


class TestPptxToText:
    def test_chart_is_omitted(self) -> None:
        # Precondition
        pptx_file = _make_pptx_with_chart()

        # Under test
        result = pptx_to_text(pptx_file)

        # Postcondition
        assert "Introduction" in result
        assert "first slide" in result
        assert "[chart omitted]" in result
        # The actual chart data should NOT appear in the output.
        assert "Revenue" not in result
        assert "Q1" not in result

    def test_text_only_pptx(self) -> None:
        # Precondition
        pptx_file = _make_pptx_without_chart()

        # Under test
        result = pptx_to_text(pptx_file)

        # Postcondition
        assert "Hello World" in result
        assert "Some content" in result
        assert "[chart omitted]" not in result
